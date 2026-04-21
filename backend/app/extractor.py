from __future__ import annotations

from collections import Counter
from pathlib import Path

from pptx import Presentation

from .heuristics import (
    infer_color_roles,
    infer_common_alignments,
    infer_component_candidates,
    infer_layout_archetypes,
    infer_pattern_diagnostics,
    infer_spacing_rhythm,
    infer_typography_roles,
)
from .models import (
    BoundingBox,
    DeckAnalysis,
    ElementAnalysis,
    PatternInfo,
    PresentationInfo,
    SlideAnalysis,
    ThemeInfo,
)
from .utils import emu_to_inches, rgb_to_hex, to_pct


def _is_group_shape(shape: object) -> bool:
    shape_type = getattr(shape, "shape_type", None)
    if shape_type is None:
        return False
    return "GROUP" in str(shape_type).upper()


def _shape_type(shape: object) -> str:
    if getattr(shape, "has_chart", False):
        return "chart_like"
    if _is_group_shape(shape):
        return "group_like"
    if getattr(shape, "has_text_frame", False):
        return "text"
    if getattr(shape, "shape_type", None) is not None and "PICTURE" in str(shape.shape_type):
        return "image"
    return "shape"


def _extract_theme_fonts(prs: Presentation) -> list[str]:
    fonts: list[str] = []
    font_scheme = getattr(getattr(getattr(prs, "slide_master", None), "theme", None), "font_scheme", None)
    if not font_scheme:
        return fonts

    for branch in ("major_font", "minor_font"):
        branch_obj = getattr(font_scheme, branch, None)
        latin = getattr(branch_obj, "latin", None)
        typeface = getattr(latin, "typeface", None)
        if typeface and typeface not in fonts:
            fonts.append(typeface)
    return fonts


def _extract_theme_colors(prs: Presentation) -> dict[str, str]:
    color_map: dict[str, str] = {}
    color_scheme = getattr(getattr(getattr(prs, "slide_master", None), "theme", None), "color_scheme", None)
    if not color_scheme:
        return color_map

    for name in (
        "dk1",
        "lt1",
        "dk2",
        "lt2",
        "accent1",
        "accent2",
        "accent3",
        "accent4",
        "accent5",
        "accent6",
        "hlink",
        "fol_hlink",
    ):
        color_obj = getattr(color_scheme, name, None)
        try:
            rgb = rgb_to_hex(getattr(color_obj, "rgb", None))
        except (AttributeError, TypeError, ValueError):
            rgb = None
        if rgb:
            color_map[name.upper()] = rgb
    return color_map


def _resolve_font_color(font: object, theme_colors: dict[str, str]) -> str | None:
    color = getattr(font, "color", None)
    if color is None:
        return None
    rgb = rgb_to_hex(getattr(color, "rgb", None))
    if rgb:
        return rgb
    theme_key = str(getattr(color, "theme_color", "")).split(".")[-1].upper()
    return theme_colors.get(theme_key)


def _text_style(
    shape: object,
    theme_fonts: list[str],
    theme_colors: dict[str, str],
) -> tuple[str | None, str | None, float | None, bool | None, bool | None, str | None]:
    if not getattr(shape, "has_text_frame", False):
        return (None, None, None, None, None, None)

    text_sample = (shape.text or "").strip().replace("\n", " ")[:180] or None
    font_name = None
    font_size = None
    bold = None
    italic = None
    text_color = None

    for paragraph in shape.text_frame.paragraphs:
        p_font = getattr(paragraph, "font", None)
        if p_font is not None:
            font_name = font_name or getattr(p_font, "name", None)
            size = getattr(p_font, "size", None)
            if font_size is None and size is not None:
                font_size = round(size.pt, 2)
            if bold is None and getattr(p_font, "bold", None) is not None:
                bold = bool(p_font.bold)
            if italic is None and getattr(p_font, "italic", None) is not None:
                italic = bool(p_font.italic)
            text_color = text_color or _resolve_font_color(p_font, theme_colors)

        for run in paragraph.runs:
            if run.text and run.text.strip():
                if run.font is not None:
                    font_name = font_name or getattr(run.font, "name", None)
                    size = getattr(run.font, "size", None)
                    if font_size is None and size is not None:
                        font_size = round(size.pt, 2)
                    if bold is None and getattr(run.font, "bold", None) is not None:
                        bold = bool(run.font.bold)
                    if italic is None and getattr(run.font, "italic", None) is not None:
                        italic = bool(run.font.italic)
                    text_color = text_color or _resolve_font_color(run.font, theme_colors)
                if font_name or font_size or text_color:
                    return (text_sample, font_name, font_size, bold, italic, text_color)

    if not font_name and theme_fonts:
        font_name = theme_fonts[0]
    return (text_sample, font_name, font_size, bold, italic, text_color)


def _shape_color(shape: object, theme_colors: dict[str, str]) -> tuple[str | None, str | None]:
    fill_color = None
    line_color = None

    if getattr(shape, "fill", None) is not None:
        try:
            fore_color = getattr(shape.fill, "fore_color", None)
            fill_color = rgb_to_hex(getattr(fore_color, "rgb", None))
            if not fill_color:
                theme_key = str(getattr(fore_color, "theme_color", "")).split(".")[-1].upper()
                fill_color = theme_colors.get(theme_key)
        except (AttributeError, TypeError, ValueError):
            fill_color = None
    if getattr(shape, "line", None) is not None:
        try:
            line = getattr(shape.line, "color", None)
            line_color = rgb_to_hex(getattr(line, "rgb", None))
            if not line_color:
                theme_key = str(getattr(line, "theme_color", "")).split(".")[-1].upper()
                line_color = theme_colors.get(theme_key)
        except (AttributeError, TypeError, ValueError):
            line_color = None

    return fill_color, line_color


def _slide_background(slide: object, theme_colors: dict[str, str]) -> str | None:
    fill = getattr(slide.background, "fill", None)
    if fill is None:
        return None
    try:
        fore = getattr(fill, "fore_color", None)
        rgb = rgb_to_hex(getattr(fore, "rgb", None))
        if rgb:
            return rgb
        theme_key = str(getattr(fore, "theme_color", "")).split(".")[-1].upper()
        return theme_colors.get(theme_key)
    except (AttributeError, TypeError, ValueError):
        return None


def _iter_shapes(shapes: object, x_offset_emu: int = 0, y_offset_emu: int = 0):
    for shape in shapes:
        left = int(getattr(shape, "left", 0) or 0) + x_offset_emu
        top = int(getattr(shape, "top", 0) or 0) + y_offset_emu
        yield shape, left, top

        if _is_group_shape(shape):
            child_shapes = getattr(shape, "shapes", None)
            if child_shapes is not None:
                yield from _iter_shapes(child_shapes, left, top)


def extract_pptx(path: str | Path) -> DeckAnalysis:
    """Extract a heuristic design model from a PPTX file."""
    prs = Presentation(str(path))
    slide_width = emu_to_inches(prs.slide_width)
    slide_height = emu_to_inches(prs.slide_height)
    theme_fonts = _extract_theme_fonts(prs)
    theme_colors = _extract_theme_colors(prs)

    presentation = PresentationInfo(
        slide_width_in=slide_width,
        slide_height_in=slide_height,
        aspect_ratio=round(slide_width / slide_height, 3) if slide_height else 0.0,
    )

    slides: list[SlideAnalysis] = []
    font_counter: Counter[str] = Counter()
    color_counter: Counter[str] = Counter()
    size_counter: Counter[float] = Counter()

    for idx, slide in enumerate(prs.slides, start=1):
        elements: list[ElementAnalysis] = []
        for shape, left_emu, top_emu in _iter_shapes(slide.shapes):
            x_in = emu_to_inches(left_emu)
            y_in = emu_to_inches(top_emu)
            w_in = emu_to_inches(getattr(shape, "width", 0))
            h_in = emu_to_inches(getattr(shape, "height", 0))

            bbox = BoundingBox(
                x_in=x_in,
                y_in=y_in,
                w_in=w_in,
                h_in=h_in,
                x_pct=to_pct(x_in, slide_width),
                y_pct=to_pct(y_in, slide_height),
                w_pct=to_pct(w_in, slide_width),
                h_pct=to_pct(h_in, slide_height),
            )

            text_sample, font_family, font_size_pt, font_bold, font_italic, text_color = _text_style(
                shape, theme_fonts, theme_colors
            )
            fill_color, line_color = _shape_color(shape, theme_colors)
            element_type = _shape_type(shape)

            element = ElementAnalysis(
                type=element_type,
                bbox=bbox,
                text_sample=text_sample,
                font_family=font_family,
                font_size_pt=font_size_pt,
                font_bold=font_bold,
                font_italic=font_italic,
                fill_color=fill_color,
                line_color=line_color,
                text_color=text_color,
            )
            elements.append(element)

            if font_family:
                font_counter[font_family] += 1
            if font_size_pt:
                size_counter[font_size_pt] += 1
            if text_color:
                color_counter[text_color] += 2
            if fill_color:
                color_counter[fill_color] += max(1, int((w_in * h_in) / 4))
            if line_color:
                color_counter[line_color] += 1

        slides.append(
            SlideAnalysis(
                slide_index=idx,
                background=_slide_background(slide, theme_colors),
                elements=elements,
            )
        )

    ranked_colors = [color for color, _ in color_counter.most_common(12)]
    theme = ThemeInfo(
        candidate_fonts=[font for font, _ in font_counter.most_common(8)],
        candidate_colors=ranked_colors,
        common_font_sizes_pt=[size for size, _ in size_counter.most_common(8)],
        theme_fonts=theme_fonts,
        theme_colors=list(theme_colors.values())[:12],
    )

    dominant_layouts = infer_layout_archetypes(slides)
    typography_roles = infer_typography_roles(slides, slide_height)
    components = infer_component_candidates(slides)
    patterns = PatternInfo(
        dominant_layouts=dominant_layouts,
        typography_roles=typography_roles,
        component_candidates=components,
        spacing_rhythm=infer_spacing_rhythm(slides, slide_width, slide_height),
        common_alignments=infer_common_alignments(slides),
        color_roles=infer_color_roles(ranked_colors),
        diagnostics=infer_pattern_diagnostics(
            slides=slides,
            typography_roles=typography_roles,
            dominant_layouts=dominant_layouts,
            component_candidates=components,
            color_count=len(ranked_colors),
        ),
    )

    return DeckAnalysis(presentation=presentation, theme=theme, slides=slides, patterns=patterns)
