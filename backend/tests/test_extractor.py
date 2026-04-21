from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from app.extractor import extract_pptx


def _make_sample_pptx(path: Path) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    title_box = slide.shapes.add_textbox(prs.slide_width // 10, prs.slide_height // 10, prs.slide_width // 2, prs.slide_height // 8)
    p = title_box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "Q2 Revenue Overview"
    run.font.name = "Calibri"
    run.font.size = 381000  # 30pt

    body_box = slide.shapes.add_textbox(prs.slide_width // 10, prs.slide_height // 3, prs.slide_width // 2, prs.slide_height // 4)
    p2 = body_box.text_frame.paragraphs[0]
    r2 = p2.add_run()
    r2.text = "Key points and growth drivers"
    r2.font.name = "Calibri"
    r2.font.size = 228600  # 18pt

    prs.save(path)


def test_extract_pptx_smoke(tmp_path: Path) -> None:
    pptx_path = tmp_path / "sample.pptx"
    _make_sample_pptx(pptx_path)

    analysis = extract_pptx(pptx_path)

    assert analysis.presentation.slide_width_in > 0
    assert len(analysis.slides) == 1
    assert any(e.type == "text" for e in analysis.slides[0].elements)
    assert analysis.theme.candidate_fonts
    assert "overall" in analysis.patterns.diagnostics.scores
