from __future__ import annotations

import io

from fastapi.testclient import TestClient
from pptx import Presentation

from app.main import app


def _pptx_bytes() -> bytes:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    box = slide.shapes.add_textbox(prs.slide_width // 8, prs.slide_height // 8, prs.slide_width // 2, prs.slide_height // 8)
    box.text_frame.text = "Executive Summary"

    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


def test_health() -> None:
    client = TestClient(app)
    res = client.get("/health")
    assert res.status_code == 200
    assert res.json()["status"] == "ok"
    assert res.headers["x-content-type-options"] == "nosniff"
    assert res.headers["x-frame-options"] == "DENY"


def test_extract_endpoint_smoke() -> None:
    client = TestClient(app)
    payload = _pptx_bytes()

    res = client.post(
        "/extract",
        files={"file": ("sample.pptx", payload, "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
    )

    assert res.status_code == 200
    body = res.json()
    assert "design_md" in body
    assert "analysis" in body
    assert "presentation" in body["analysis"]
    assert body["run_id"]


def test_extract_batch_and_save_smoke() -> None:
    client = TestClient(app)
    payload_one = _pptx_bytes()
    payload_two = _pptx_bytes()

    res = client.post(
        "/extract/batch",
        files=[
            ("files", ("a.pptx", payload_one, "application/vnd.openxmlformats-officedocument.presentationml.presentation")),
            ("files", ("b.pptx", payload_two, "application/vnd.openxmlformats-officedocument.presentationml.presentation")),
        ],
    )
    assert res.status_code == 200
    body = res.json()
    assert len(body["results"]) == 2
    run_id = body["results"][0]["run_id"]

    save_res = client.post(
        "/design/save",
        json={"run_id": run_id, "design_md": "# DESIGN.md\n\nEdited"},
    )
    assert save_res.status_code == 200
    save_body = save_res.json()
    assert save_body["run_id"] == run_id
    assert save_body["saved_path"].endswith("design.edited.md")


def test_extract_rejects_non_pptx_content() -> None:
    client = TestClient(app)
    res = client.post(
        "/extract",
        files={"file": ("evil.pptx", b"not-a-zip", "application/octet-stream")},
    )
    assert res.status_code == 400
    assert "Invalid PPTX" in res.json()["detail"]


def test_save_design_rejects_invalid_run_id() -> None:
    client = TestClient(app)
    res = client.post(
        "/design/save",
        json={"run_id": "../../tmp/owned", "design_md": "# DESIGN.md\n\nx"},
    )
    assert res.status_code == 404
